import os
import logging
import traceback
from urllib.parse import urlparse
import io
from PIL import Image  # Add this import
from pdf2image import convert_from_path  # For PDF files
from pptx import Presentation  # For PPTX files

from app.utils.deterministic_checks import format_check, size_check, slide_count_check
from app.utils.file_analyzers import (analyze_pptx, analyze_pdf,
                                      analyze_markdown, analyze_keynote,
                                      analyze_google_slides, analyze_canva,
                                      analyze_figma)
from app.utils.probabilistic_checks import analyze_slide_image
from app.models.schemas import (FileAnalysisResult, DeterministicCheckResult,
                                ProbabilisticCheckResult, SlideAnalysis,
                                AnalysisResponse, TitleSlideCheck, ImageCheck,
                                BulletPointCheck)
from pydantic import ValidationError

# Get the logger
logger = logging.getLogger("slide_analyzer")


def is_url(path):
    try:
        result = urlparse(path)
        return all([result.scheme, result.netloc])
    except ValueError:
        return False


async def process_slide_deck(input_path: str,
                             processing_id: str,
                             deck_format: str) -> AnalysisResponse:
    logger.info(f"Starting to process slide deck: {input_path}")
    try:
        # Step 1: Format Identification
        is_file = os.path.isfile(input_path)
        is_url_path = is_url(input_path)

        logger.debug(f"Is file: {is_file}, Is URL: {is_url_path}")

        if not is_file and not is_url_path:
            logger.error(f"Invalid input: {input_path}")
            return None

        # Determine format
        if is_file:
            format_result = format_check(input_path)
            file_format = format_result['file_type']
            logger.info(f"File format determined: {file_format}")
        else:
            # For URLs, we need to determine the format based on the URL structure
            if 'docs.google.com' in input_path:
                file_format = 'google_slides'
            elif 'figma.com' in input_path:
                file_format = 'figma'
            elif 'canva.com' in input_path:
                file_format = 'canva'
            else:
                logger.error(f"Unsupported URL format: {input_path}")
                return None
            logger.info(f"URL format determined: {file_format}")

        # Step 2: Deterministic Checks
        logger.info("Performing deterministic checks")
        deterministic_checks = DeterministicCheckResult(
            format_check=format_check(input_path) if is_file else {
                "accepted_format": True,
                "file_type": file_format,
                "message": "URL format accepted"
            },
            size_check=size_check(
                os.path.getsize(input_path) / (1024 * 1024))
            if is_file else {
                "size_within_limit": True,
                "file_size_mb": 0,
                "message": "Size check not applicable for URLs"
            },
            slide_count_check={
                "slide_count_within_limit": True,
                "slide_count": 0,
                "message":
                "Slide count will be determined in detailed analysis"
            })
        logger.debug(f"Deterministic checks result: {deterministic_checks}")

        # Step 3: Detailed File Analysis and Slide Image Generation
        logger.info(
            f"Starting detailed file analysis for format: {file_format}")

        slide_images = []

        if file_format in ['pdf', '.pdf']:
            # Extract slides from PDF
            slides = convert_from_path(input_path)
            total_slides = len(slides)

            for i, slide in enumerate(slides):
                slide_number = i + 1
                image_buffer = io.BytesIO()
                slide.save(image_buffer, format='PNG')
                slide_images.append(image_buffer.getvalue())
                logger.debug(f"Processed slide image {slide_number}")

        elif file_format in ['pptx', '.pptx']:
            # Extract slides from PPTX
            prs = Presentation(input_path)
            total_slides = len(prs.slides)

            for i, slide in enumerate(prs.slides):
                slide_number = i + 1
                image_path = get_slide_image_path(input_path, slide_number)
                # Convert slide to image (requires additional tools like `python-pptx` and `Pillow`)
                # This is a placeholder for actual slide image extraction logic
                # You might need to use `pptx2png` or other methods to convert slides to images
                # Example using `slide.shapes` and `Pillow` to render image
                # Since `python-pptx` does not support slide rendering, consider using `comtypes` on Windows
                logger.warning(
                    "Slide image extraction for PPTX is not implemented.")
                # Save placeholder image
                Image.new('RGB', (1280, 720),
                          color='gray').save(image_path)
                logger.debug(
                    f"Saved placeholder slide image: {image_path}")

        else:
            logger.error(
                f"Unsupported file format for slide image extraction: {
                    file_format}"
            )
            return None

        # Step 4: Probabilistic Checks using GPT-4 on slide images
        logger.info("Starting probabilistic checks")
        probabilistic_checks = []
        title_slide_present = False
        total_bullet_points = 0
        total_images = 0

        for index, image_data in enumerate(slide_images):
            slide_number = index + 1
            logger.debug(f"Analyzing slide image {slide_number}")

            # Call your analysis function
            analysis = await analyze_slide_image(image_data, slide_number)

            # Check if analysis contains required keys
            required_keys = [
                'is_title_slide', 'bullet_points', 'images',
                'adheres_to_best_practices', 'suggestions'
            ]
            if all(key in analysis for key in required_keys):
                probabilistic_checks.append(
                    SlideAnalysis(slide_number=slide_number,
                                  analysis=analysis))

                # Aggregate data for checks
                if analysis.get('is_title_slide'):
                    title_slide_present = True
                total_bullet_points += analysis.get('bullet_points', 0)
                total_images += analysis.get('images', 0)
            else:
                logging.error(
                    f"Analysis for slide {
                        slide_number} is incomplete. Skipping."
                )

            # When saving the slide images, use the processing_id
            image_path = get_slide_image_path(processing_id, slide_number)
            with open(image_path, 'wb') as f:
                f.write(image_data)
            logger.debug(f"Saved slide image: {image_path}")

        # Create the required checks
        title_slide_check = TitleSlideCheck(
            has_title_slide=title_slide_present,
            message="Title slide is present."
            if title_slide_present else "Title slide is missing.")

        max_bullet_points = 10  # Adjust as needed
        bullet_point_check = BulletPointCheck(
            has_few_bullet_points=total_bullet_points <= max_bullet_points,
            message=f"Total bullet points: {total_bullet_points}.")

        image_check = ImageCheck(has_images=total_images > 0,
                                 image_count=total_images,
                                 message=f"Total images: {total_images}.")

        probabilistic_checks_result = ProbabilisticCheckResult(
            title_slide_check=title_slide_check,
            bullet_point_check=bullet_point_check,
            image_check=image_check,
            slide_analyses=probabilistic_checks)

        # Step 5: Construct and validate the final response
        logger.info("Constructing final response")
        analysis_response = AnalysisResponse(
            processing_id=processing_id,
            deterministic_checks=deterministic_checks,
            file_analysis=FileAnalysisResult(
                number_of_slides=total_slides,
                fonts_used=[],  # Placeholder, replace with actual data
                video_present=False,  # Placeholder, replace with actual data
                audio_present=False  # Placeholder, replace with actual data
            ),
            probabilistic_checks=probabilistic_checks_result
            # Provide other required fields or make them optional
        )
        logger.info("Analysis response created successfully")
        return analysis_response

    except Exception as e:
        logger.error(f"Error in process_slide_deck: {str(e)}")
        logger.error(traceback.format_exc())
        return None


def get_slide_image_path(processing_id: str, slide_number: int) -> str:
    """
    Generate the path for a slide image based on the processing ID and slide number.
    """
    # Create a directory for the slides if it doesn't exist
    slides_dir = os.path.join('uploads', processing_id)
    os.makedirs(slides_dir, exist_ok=True)

    # Generate the path for the specific slide image
    return os.path.join(slides_dir, f"slide_{slide_number}.png")
