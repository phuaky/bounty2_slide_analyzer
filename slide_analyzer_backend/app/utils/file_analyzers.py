import os
import zipfile
import re
import requests
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import fitz  # PyMuPDF
from google.oauth2.credentials import Credentials
from googleapiclient.discovery import build


def analyze_pptx(file_path):
    prs = Presentation(file_path)
    num_slides = len(prs.slides)
    fonts_used = set()
    video_present = False
    audio_present = False

    for slide in prs.slides:
        for shape in slide.shapes:
            # Extract fonts
            if hasattr(shape, "text_frame"):
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        font = run.font
                        if font.name:
                            fonts_used.add(font.name)
            # Check for media
            if shape.shape_type == MSO_SHAPE_TYPE.MEDIA:
                if hasattr(shape, 'media_format'):
                    media_format = shape.media_format
                    if media_format.type == 'video':
                        video_present = True
                    elif media_format.type == 'audio':
                        audio_present = True
                else:
                    video_present = True  # Cannot distinguish, assume video

    return {
        'number_of_slides': num_slides,
        'fonts_used': list(fonts_used),
        'video_present': video_present,
        'audio_present': audio_present,
    }


def analyze_pdf(file_path):
    try:
        doc = fitz.open(file_path)
    except Exception as e:
        return {'error': f'Failed to open PDF: {str(e)}'}

    num_pages = doc.page_count  # Number of slides/pages
    fonts_used = set()
    video_present = False
    audio_present = False
    slide_images = []

    # Get fonts used and render slide images
    for page in doc:
        # Extract fonts
        font_list = page.get_fonts()
        for font in font_list:
            font_name = font[3]  # Font name
            # Clean the font name by stripping the prefix (if present)
            if "+" in font_name:
                font_name = font_name.split("+")[1]
            fonts_used.add(font_name)

        # Check for annotations (embedded media)
        annots = page.annots()
        if annots:
            for annot in annots:
                subtype = annot.info['Subtype']
                if subtype == 'RichMedia':
                    video_present = True
                elif subtype == 'Sound':
                    audio_present = True

        # Render page to image
        pix = page.get_pixmap()
        image_bytes = pix.tobytes("png")
        slide_images.append(image_bytes)

    return {
        'number_of_slides': num_pages,
        'fonts_used': list(fonts_used),
        'video_present': video_present,
        'audio_present': audio_present,
        'slide_images': slide_images,  # Include slide images
    }


def analyze_markdown(file_path):
    with open(file_path, 'r', encoding='utf-8') as f:
        content = f.read()

    # Count the number of slides (assuming '---' or '***' as slide separators)
    slides = re.split(r'\n---\n|\n\*\*\*\n', content)
    num_slides = len(slides)

    # Fonts are generally not specified in Markdown
    fonts_used = []

    # Check for embedded video and audio
    video_present = bool(
        re.search(r'!\[.*\]\(.*\.(mp4|avi|mov|wmv)\)', content, re.IGNORECASE)
        or re.search(r'<video.*?>', content, re.IGNORECASE))
    audio_present = bool(re.search(r'<audio.*?>', content, re.IGNORECASE))

    return {
        'number_of_slides': num_slides,
        'fonts_used': fonts_used,
        'video_present': video_present,
        'audio_present': audio_present,
    }


def analyze_keynote(file_path):
    # Attempt to unzip the Keynote file
    try:
        with zipfile.ZipFile(file_path, 'r') as z:
            # Check for index.apxl or other XML files
            names = z.namelist()
            if 'index.apxl' in names:
                index_data = z.read('index.apxl')
                # Parse the XML data (complex and may not be reliable)
                num_slides = index_data.count(b'<slide-uuid>')
                fonts_used = []  # Not easily extractable
                video_present = False
                audio_present = False
                return {
                    'number_of_slides': num_slides,
                    'fonts_used': fonts_used,
                    'video_present': video_present,
                    'audio_present': audio_present,
                }
            else:
                return {'error': 'Unsupported Keynote file structure.'}
    except zipfile.BadZipFile:
        return {'error': 'File is not a valid Keynote package.'}


def analyze_google_slides(presentation_url):
    # Extract presentation ID from URL
    m = re.match(r'https://docs.google.com/presentation/d/([a-zA-Z0-9-_]+)',
                 presentation_url)
    if not m:
        return {'error': 'Invalid Google Slides URL'}

    presentation_id = m.group(1)

    # Load credentials from environment variables or a file
    creds = Credentials(
        token='YOUR_ACCESS_TOKEN')  # Replace with your access token

    service = build('slides', 'v1', credentials=creds)

    # Get the presentation
    presentation = service.presentations().get(
        presentationId=presentation_id).execute()

    num_slides = len(presentation.get('slides', []))

    # Fonts used
    fonts_used = set()
    video_present = False
    audio_present = False

    for slide in presentation.get('slides', []):
        for element in slide.get('pageElements', []):
            if 'shape' in element and 'text' in element['shape']:
                text_elements = element['shape']['text'].get(
                    'textElements', [])
                for te in text_elements:
                    if 'textRun' in te and 'style' in te['textRun']:
                        font_family = te['textRun']['style'].get('fontFamily')
                        if font_family:
                            fonts_used.add(font_family)
            # Check for video (videos are not directly accessible via API)
            if 'video' in element:
                video_present = True

    return {
        'number_of_slides': num_slides,
        'fonts_used': list(fonts_used),
        'video_present': video_present,
        'audio_present': audio_present,  # Not accessible via API
    }


# Example usage
# result = analyze_google_slides('https://docs.google.com/presentation/d/your_presentation_id')
# print(result)


def analyze_canva(url):
    return {
        'error':
        'Canva does not provide an API for analysis. Please export your slides as PDF or PPTX.'
    }


def analyze_figma(figma_url):
    # Replace 'YOUR_FIGMA_ACCESS_TOKEN' with your actual token
    file_key_match = re.search(r'/design/([^/]+)', figma_url)
    if not file_key_match:
        return 'Invalid Figma URL format'
    
    file_key = file_key_match.group(1)
    api_url = f'https://api.figma.com/v1/files/{file_key}'
    headers = {'X-Figma-Token': 'YOUR_FIGMA_ACCESS_TOKEN'}  # Replace with your Figma API token
    
    response = requests.get(api_url, headers=headers)
    if response.status_code != 200:
        return f'Error fetching data from Figma API: {response.status_code}'

    data = response.json()

    slides_count = 0
    fonts = set()
    video_media = False
    video_media_count = 0  # New variable to count video media instances

    contains_audio = False

    def traverse_node(node, level=0):
        nonlocal slides_count, video_media, contains_audio, video_media_count
        
        indent = "  " * level
        #print(f"{indent}Node: {node.get('name', 'Unnamed')} (Type: {node['type']})")
        
        if node['type'] == 'CANVAS':
            slides_count = len(node.get('children', []))
            print(f"{indent}Found {slides_count} slides in {node.get('name', 'Unnamed')}")

        if'style' in node and 'fontFamily' in node['style']:
            fonts.add(node['style']['fontFamily'])
            #print(f"{indent}Font found: {node['style']['fontFamily']}")
            
            # **Check for Video Fill in Nodes**
        if 'interactions' in node:
            for interaction in node['interactions']:
                if 'actions' in interaction:
                    for action in interaction['actions']:
                        if action.get('type') == 'UPDATE_MEDIA_RUNTIME' and action.get('mediaAction') == 'TOGGLE_PLAY_PAUSE':
                            video_media = True
                            video_media_count += 1  # Increment video media count
                            print(f"{indent}Video media found in {node.get('name', 'Unnamed')}")
                                                                           
        if 'children' in node:
            # print(f"{indent}Traversing {len(node['children'])} children of {node.get('name', 'Unnamed')}")
            for child in node['children']:
                traverse_node(child, level + 1)

    print("Starting document traversal")
    traverse_node(data['document'], 0)

    print(f"\nTotal slides found: {slides_count}")
    print(f"Fonts found: {fonts}")
    print(f"Video media present: {video_media}")
    print(f"Video media count: {video_media_count}")
    print(f"Audio in video: {contains_audio}")

    return {
        'slides_count': slides_count,
        'fonts': list(fonts),
        'video_media': video_media,
        'video_media_count': video_media_count,
        'contains_audio': contains_audio
    }
