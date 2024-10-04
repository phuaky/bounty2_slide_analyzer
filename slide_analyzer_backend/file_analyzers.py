import os
import zipfile
import re
import requests
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import fitz  # PyMuPDF
from google.oauth2.credentials import Credentials
from googleapiclient.discovery import build


def analyze_pptx(file_path):
    prs = Presentation(file_path)
    num_slides = len(prs.slides)
    fonts_used = set()
    video_present = False
    audio_present = False

    for slide in prs.slides:
        for shape in slide.shapes:
            # Extract fonts
            if hasattr(shape, "text_frame"):
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        font = run.font
                        if font.name:
                            fonts_used.add(font.name)
            # Check for media
            if shape.shape_type == MSO_SHAPE_TYPE.MEDIA:
                if hasattr(shape, 'media_format'):
                    media_format = shape.media_format
                    if media_format.type == 'video':
                        video_present = True
                    elif media_format.type == 'audio':
                        audio_present = True
                else:
                    video_present = True  # Cannot distinguish, assume video

    return {
        'number_of_slides': num_slides,
        'fonts_used': list(fonts_used),
        'video_present': video_present,
        'audio_present': audio_present,
    }


def analyze_pdf(file_path):
    doc = fitz.open(file_path)
    num_pages = doc.page_count  # Number of slides/pages
    fonts_used = set()
    video_present = False
    audio_present = False

    # Get fonts used
    for page in doc:
        font_list = page.get_fonts()
        for font in font_list:
            font_name = font[3]  # Font name
            # Clean the font name by stripping the prefix (if present)
            if "+" in font_name:
                font_name = font_name.split("+")[1]
            fonts_used.add(font_name)

        # Check for annotations (embedded media)
        annots = page.annots()
        if annots:
            for annot in annots:
                subtype = annot.info['Subtype']
                if subtype == 'RichMedia':
                    video_present = True
                elif subtype == 'Sound':
                    audio_present = True

    return {
        'number_of_slides': num_pages,
        'fonts_used': list(fonts_used),
        'video_present': video_present,
        'audio_present': audio_present,
    }


def analyze_markdown(file_path):
    with open(file_path, 'r', encoding='utf-8') as f:
        content = f.read()

    # Count the number of slides (assuming '---' or '***' as slide separators)
    slides = re.split(r'\n---\n|\n\*\*\*\n', content)
    num_slides = len(slides)

    # Fonts are generally not specified in Markdown
    fonts_used = []

    # Check for embedded video and audio
    video_present = bool(
        re.search(r'!\[.*\]\(.*\.(mp4|avi|mov|wmv)\)', content, re.IGNORECASE)
        or re.search(r'<video.*?>', content, re.IGNORECASE))
    audio_present = bool(re.search(r'<audio.*?>', content, re.IGNORECASE))

    return {
        'number_of_slides': num_slides,
        'fonts_used': fonts_used,
        'video_present': video_present,
        'audio_present': audio_present,
    }


def analyze_keynote(file_path):
    # Attempt to unzip the Keynote file
    try:
        with zipfile.ZipFile(file_path, 'r') as z:
            # Check for index.apxl or other XML files
            names = z.namelist()
            if 'index.apxl' in names:
                index_data = z.read('index.apxl')
                # Parse the XML data (complex and may not be reliable)
                num_slides = index_data.count(b'<slide-uuid>')
                fonts_used = []  # Not easily extractable
                video_present = False
                audio_present = False
                return {
                    'number_of_slides': num_slides,
                    'fonts_used': fonts_used,
                    'video_present': video_present,
                    'audio_present': audio_present,
                }
            else:
                return {'error': 'Unsupported Keynote file structure.'}
    except zipfile.BadZipFile:
        return {'error': 'File is not a valid Keynote package.'}


def analyze_google_slides(presentation_url):
    # Extract presentation ID from URL
    m = re.match(r'https://docs.google.com/presentation/d/([a-zA-Z0-9-_]+)',
                 presentation_url)
    if not m:
        return {'error': 'Invalid Google Slides URL'}

    presentation_id = m.group(1)

    # Load credentials from environment variables or a file
    creds = Credentials(
        # Fetch from environment variables
        token=os.getenv('GOOGLE_ACCESS_TOKEN')
    )

    service = build('slides', 'v1', credentials=creds)

    # Get the presentation
    presentation = service.presentations().get(
        presentationId=presentation_id).execute()

    num_slides = len(presentation.get('slides', []))

    # Fonts used
    fonts_used = set()
    video_present = False
    audio_present = False

    for slide in presentation.get('slides', []):
        for element in slide.get('pageElements', []):
            if 'shape' in element and 'text' in element['shape']:
                text_elements = element['shape']['text'].get(
                    'textElements', [])
                for te in text_elements:
                    if 'textRun' in te and 'style' in te['textRun']:
                        font_family = te['textRun']['style'].get('fontFamily')
                        if font_family:
                            fonts_used.add(font_family)
            # Check for video (videos are not directly accessible via API)
            if 'video' in element:
                video_present = True

    return {
        'number_of_slides': num_slides,
        'fonts_used': list(fonts_used),
        'video_present': video_present,
        'audio_present': audio_present,  # Not accessible via API
    }


# Example usage
# result = analyze_google_slides('https://docs.google.com/presentation/d/your_presentation_id')
# print(result)


def analyze_canva(url):
    return {
        'error':
        'Canva does not provide an API for analysis. Please export your slides as PDF or PPTX.'
    }


def analyze_figma(file_url):
    # Replace 'YOUR_FIGMA_ACCESS_TOKEN' with your actual token
    access_token = os.getenv(
        'FIGMA_ACCESS_TOKEN')  # Store your token in Replit secrets
    if not access_token:
        return {
            'error':
            'Figma access token not found. Please set FIGMA_ACCESS_TOKEN in environment variables.'
        }

    # Extract file key from URL
    m = re.search(r'file/([a-zA-Z0-9]+)', file_url)
    if not m:
        return {'error': 'Invalid Figma file URL'}

    file_key = m.group(1)

    # Get the file data
    url = f'https://api.figma.com/v1/files/{file_key}'
    response = requests.get(url, headers=headers)
    if response.status_code != 200:
        return {'error': f'Error fetching Figma file: {response.text}'}

    data = response.json()

    num_slides = 0
    fonts_used = set()
    video_present = False
    audio_present = False

    def traverse(node):
        nonlocal num_slides, fonts_used
        if node['type'] == 'FRAME':
            num_slides += 1
        if 'style' in node:
            font_family = node['style'].get('fontFamily')
            if font_family:
                fonts_used.add(font_family)
        if 'children' in node:
            for child in node['children']:
                traverse(child)

    document = data['document']
    traverse(document)

    return {
        'number_of_slides': num_slides,
        'fonts_used': list(fonts_used),
        'video_present': video_present,  # Not directly accessible via API
        'audio_present': audio_present,  # Not directly accessible via API
    }
